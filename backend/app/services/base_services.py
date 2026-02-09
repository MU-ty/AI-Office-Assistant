"""
其他模块的Service框架 - 占位符
这些模块的完整实现将逐个在后续开发中完成
"""

from typing import Optional, List, Dict, Any
import json
from datetime import datetime
import os
import uuid
import asyncio
import aiohttp
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy import select, update, delete, and_
from sqlalchemy.orm import selectinload

from app.models.polish import PolishTask, PolishIssue
from app.models.translation import TranslationTask, TranslationTerminology
from app.models.ppt import PPTProject
from app.core.config import settings
from app.services.polish_normalization_service import AcademicNormalizationService
from app.services.llm_service import llm_service
from app.utils.logger import get_logger

logger = get_logger(__name__)


# ============================================================
# 学术润色服务
# ============================================================

class PolishService:
    """学术润色服务 - 实现流程图中的学术规范化模块"""
    
    def __init__(self, db: AsyncSession):
        self.db = db
        self.normalization_service = AcademicNormalizationService()
    
    async def create_task(self, task_data: Dict, user_id: int) -> dict:
        """
        创建润色任务
        
        流程：
        1. 验证输入
        2. 创建任务记录
        3. 调用Qwen LLM进行深度学术化改写
        4. 执行学术规范化分析作为辅助检查
        5. 保存问题到数据库
        6. 返回任务结果
        """
        logger.info("创建学术润色任务")
        
        try:
            original_text = task_data.get("original_text", "")
            polish_level = task_data.get("polish_level", "standard")
            auto_fix_enabled = task_data.get("auto_fix_enabled", False)
            document_id = task_data.get("document_id")
            
            if not original_text or len(original_text) < 1:
                raise ValueError("文本不能为空")
            
            # 创建任务
            task = PolishTask(
                user_id=user_id,
                original_text=original_text,
                polish_level=polish_level,
                auto_fix_enabled=str(auto_fix_enabled).lower(),
                document_id=document_id,
                status="processing"
            )
            self.db.add(task)
            await self.db.flush()  # 获取task ID
            
            # 首先调用Qwen LLM进行深度学术化改写
            logger.info(f"调用Qwen API进行学术化改写，任务ID: {task.id}")
            polished_text = await self._polish_with_llm(original_text, polish_level)
            task.polished_text = polished_text
            
            # 然后执行规则引擎分析作为辅助检查
            logger.info("执行规则引擎分析")
            analysis_result = self.normalization_service.analyze_text(polished_text)
            
            # 保存问题到数据库
            all_issues = []
            issue_type_map = {
                "terminology_issues": "terminology",
                "tense_issues": "tense",
                "style_issues": "style",
                "thesis_issues": "thesis"
            }
            
            for key, issue_type in issue_type_map.items():
                for issue_data in analysis_result.get(key, []):
                    polish_issue = PolishIssue(
                        task_id=task.id,
                        issue_type=issue_type,
                        severity=issue_data.get("severity", "medium"),
                        location={
                            "start": issue_data["start"],
                            "end": issue_data["end"]
                        },
                        original_content=issue_data["original"],
                        suggested_content=issue_data["suggested"],
                        reason=issue_data.get("reason", ""),
                        rule_id=issue_data.get("rule_id"),
                        confidence=issue_data.get("confidence", 0.0)
                    )
                    self.db.add(polish_issue)
                    all_issues.append(polish_issue)
            
            # 更新任务统计
            task.total_issues = analysis_result["total_issues"]
            task.terminology_issues = analysis_result.get("terminology_issues", [])
            task.tense_issues = analysis_result.get("tense_issues", [])
            task.style_issues = analysis_result.get("style_issues", [])
            task.thesis_issues = analysis_result.get("thesis_issues", [])
            
            # 如果启用自动修复，应用修复
            if auto_fix_enabled:
                polished_text, fixed_count = self.normalization_service.apply_fixes(
                    polished_text,
                    all_issues
                )
                task.polished_text = polished_text
                task.fixed_issues = fixed_count
                
                # 计算准确率
                if task.total_issues > 0:
                    task.accuracy = fixed_count / task.total_issues
            else:
                task.polished_text = polished_text
                task.fixed_issues = 0
                task.accuracy = 0.0
            
            task.status = "completed"
            task.completed_at = datetime.utcnow()
            
            await self.db.commit()
            await self.db.refresh(task)
            
            logger.info(f"任务创建成功，ID: {task.id}, 问题数: {task.total_issues}")
            return self._format_task_response(task)
            
        except Exception as e:
            await self.db.rollback()
            logger.error(f"创建任务失败: {str(e)}")
            raise
    
    async def list_tasks(
        self, user_id: int, skip: int = 0, limit: int = 10, status: Optional[str] = None
    ) -> dict:
        """获取任务列表"""
        try:
            query = select(PolishTask).where(PolishTask.user_id == user_id)
            
            if status:
                query = query.where(PolishTask.status == status)
            
            query = query.offset(skip).limit(limit).order_by(PolishTask.created_at.desc())
            result = await self.db.execute(query)
            tasks = result.scalars().all()
            
            # 统计总数
            count_query = select(PolishTask).where(PolishTask.user_id == user_id)
            if status:
                count_query = count_query.where(PolishTask.status == status)
            count_result = await self.db.execute(count_query)
            total = len(count_result.scalars().all())
            
            return {
                "total": total,
                "skip": skip,
                "limit": limit,
                "items": [self._format_task_response(task) for task in tasks]
            }
        except Exception as e:
            logger.error(f"获取任务列表失败: {str(e)}")
            raise
    
    async def get_task(self, task_id: int, user_id: int) -> dict:
        """获取任务详情"""
        try:
            query = select(PolishTask).where(
                and_(PolishTask.id == task_id, PolishTask.user_id == user_id)
            )
            result = await self.db.execute(query)
            task = result.scalars().first()
            
            if not task:
                raise ValueError(f"任务 {task_id} 不存在")
            
            return self._format_task_response(task)
        except Exception as e:
            logger.error(f"获取任务详情失败: {str(e)}")
            raise
    
    async def update_task(self, task_id: int, task_data: Dict, user_id: int) -> dict:
        """更新任务"""
        try:
            query = select(PolishTask).where(
                and_(PolishTask.id == task_id, PolishTask.user_id == user_id)
            )
            result = await self.db.execute(query)
            task = result.scalar_one_or_none()
            if not task:
                raise ValueError(f"任务 {task_id} 不存在")
            
            # 更新允许的字段
            if "original_text" in task_data:
                task.original_text = task_data["original_text"]
                # 重新分析
                analysis_result = self.normalization_service.analyze_text(task.original_text)
                task.total_issues = analysis_result["total_issues"]
                task.status = "completed"
            
            if "polish_level" in task_data:
                task.polish_level = task_data["polish_level"]
            
            if "auto_fix_enabled" in task_data:
                task.auto_fix_enabled = str(task_data["auto_fix_enabled"]).lower()
            
            task.updated_at = datetime.utcnow()
            
            await self.db.commit()
            await self.db.refresh(task)
            
            logger.info(f"任务 {task_id} 已更新")
            return self._format_task_response(task)
        except Exception as e:
            await self.db.rollback()
            logger.error(f"更新任务失败: {str(e)}")
            raise
    
    async def delete_task(self, task_id: int, user_id: int) -> None:
        """删除任务"""
        try:
            query = select(PolishTask).where(
                and_(PolishTask.id == task_id, PolishTask.user_id == user_id)
            )
            result = await self.db.execute(query)
            task = result.scalar_one_or_none()
            if not task:
                raise ValueError(f"任务 {task_id} 不存在")
            
            # 删除相关的问题记录
            await self.db.execute(
                delete(PolishIssue).where(PolishIssue.task_id == task_id)
            )
            
            # 删除任务
            await self.db.delete(task)
            await self.db.commit()
            
            logger.info(f"任务 {task_id} 已删除")
        except Exception as e:
            await self.db.rollback()
            logger.error(f"删除任务失败: {str(e)}")
            raise
    
    async def get_issues(
        self, task_id: int, user_id: int, filter_type: Optional[str] = None
    ) -> dict:
        """获取问题列表"""
        try:
            query = select(PolishTask).where(
                and_(PolishTask.id == task_id, PolishTask.user_id == user_id)
            )
            result = await self.db.execute(query)
            task = result.scalar_one_or_none()
            if not task:
                raise ValueError(f"任务 {task_id} 不存在")
            
            query = select(PolishIssue).where(PolishIssue.task_id == task_id)
            
            if filter_type:
                query = query.where(PolishIssue.issue_type == filter_type)
            
            query = query.order_by(PolishIssue.location)
            result = await self.db.execute(query)
            issues = result.scalars().all()
            
            return {
                "task_id": task_id,
                "total": len(issues),
                "issues": [self._format_issue_response(issue) for issue in issues]
            }
        except Exception as e:
            logger.error(f"获取问题列表失败: {str(e)}")
            raise
    
    async def accept_suggestion(
        self,
        task_id: int,
        issue_id: int,
        feedback: Optional[str] = None,
        user_id: Optional[int] = None
    ) -> dict:
        """接受建议"""
        try:
            issue = await self.db.get(PolishIssue, issue_id)
            if not issue or issue.task_id != task_id:
                raise ValueError(f"问题 {issue_id} 不存在或不属于任务 {task_id}")

            task = await self.db.get(PolishTask, task_id)
            if not task or (user_id is not None and task.user_id != user_id):
                raise ValueError(f"任务 {task_id} 不存在")
            
            issue.status = "accepted"
            issue.accepted_at = datetime.utcnow()
            
            # 更新任务的修复计数
            if task:
                task.fixed_issues = (task.fixed_issues or 0) + 1
                if task.total_issues > 0:
                    task.accuracy = task.fixed_issues / task.total_issues
                task.polished_text = self._apply_suggestion(
                    task.polished_text or task.original_text,
                    issue.location,
                    issue.suggested_content
                )
            
            await self.db.commit()
            await self.db.refresh(issue)
            
            logger.info(f"已接受问题 {issue_id} 的建议")
            return self._format_issue_response(issue)
        except Exception as e:
            await self.db.rollback()
            logger.error(f"接受建议失败: {str(e)}")
            raise
    
    async def reject_suggestion(
        self,
        task_id: int,
        issue_id: int,
        reason: Optional[str] = None,
        user_id: Optional[int] = None
    ) -> dict:
        """拒绝建议"""
        try:
            issue = await self.db.get(PolishIssue, issue_id)
            if not issue or issue.task_id != task_id:
                raise ValueError(f"问题 {issue_id} 不存在或不属于任务 {task_id}")

            task = await self.db.get(PolishTask, task_id)
            if not task or (user_id is not None and task.user_id != user_id):
                raise ValueError(f"任务 {task_id} 不存在")
            
            issue.status = "rejected"
            
            await self.db.commit()
            await self.db.refresh(issue)
            
            logger.info(f"已拒绝问题 {issue_id} 的建议")
            return self._format_issue_response(issue)
        except Exception as e:
            await self.db.rollback()
            logger.error(f"拒绝建议失败: {str(e)}")
            raise
    
    async def export_result(self, task_id: int, format_type: str = "json", user_id: Optional[int] = None) -> dict:
        """导出结果"""
        try:
            task = await self.db.get(PolishTask, task_id)
            if not task:
                raise ValueError(f"任务 {task_id} 不存在")
            if user_id is not None and task.user_id != user_id:
                raise ValueError(f"任务 {task_id} 不存在")
            
            # 获取所有问题
            query = select(PolishIssue).where(PolishIssue.task_id == task_id)
            result = await self.db.execute(query)
            issues = result.scalars().all()
            
            if format_type == "json":
                return {
                    "task": self._format_task_response(task),
                    "issues": [self._format_issue_response(issue) for issue in issues],
                    "export_time": datetime.utcnow().isoformat()
                }
            elif format_type == "txt":
                return self._export_as_text(task, issues)
            else:
                raise ValueError(f"不支持的导出格式: {format_type}")
        except Exception as e:
            logger.error(f"导出结果失败: {str(e)}")
            raise
    
    # ================================================================
    # 辅助方法
    # ================================================================
    
    def _format_task_response(self, task: PolishTask) -> dict:
        """格式化任务响应"""
        return {
            "id": task.id,
            "user_id": task.user_id,
            "document_id": task.document_id,
            "original_text": task.original_text[:500] + "..." if len(task.original_text) > 500 else task.original_text,
            "polished_text": task.polished_text[:500] + "..." if task.polished_text and len(task.polished_text) > 500 else task.polished_text,
            "status": task.status,
            "polish_level": task.polish_level,
            "total_issues": task.total_issues,
            "fixed_issues": task.fixed_issues,
            "accuracy": round(task.accuracy, 2) if task.accuracy else 0.0,
            "auto_fix_enabled": task.auto_fix_enabled,
            "created_at": task.created_at.isoformat() if task.created_at else None,
            "updated_at": task.updated_at.isoformat() if task.updated_at else None,
            "completed_at": task.completed_at.isoformat() if task.completed_at else None,
        }
    
    def _format_issue_response(self, issue: PolishIssue) -> dict:
        """格式化问题响应"""
        return {
            "id": issue.id,
            "task_id": issue.task_id,
            "issue_type": issue.issue_type,
            "severity": issue.severity,
            "location": issue.location,
            "original_content": issue.original_content,
            "suggested_content": issue.suggested_content,
            "reason": issue.reason,
            "status": issue.status,
            "rule_id": issue.rule_id,
            "confidence": round(issue.confidence, 2) if issue.confidence else 0.0,
            "accepted_at": issue.accepted_at.isoformat() if issue.accepted_at else None,
            "created_at": issue.created_at.isoformat() if issue.created_at else None,
        }
    
    def _export_as_text(self, task: PolishTask, issues: List[PolishIssue]) -> str:
        """导出为文本格式"""
        lines = []
        lines.append("=" * 70)
        lines.append("学术润色任务导出报告")
        lines.append("=" * 70)
        lines.append(f"\n任务ID: {task.id}")
        lines.append(f"创建时间: {task.created_at}")
        lines.append(f"状态: {task.status}")
        lines.append(f"润色级别: {task.polish_level}")
        lines.append(f"\n统计信息:")
        lines.append(f"  总问题数: {task.total_issues}")
        lines.append(f"  已修复: {task.fixed_issues}")
        lines.append(f"  准确率: {task.accuracy * 100:.1f}%")
        lines.append(f"\n原始文本:")
        lines.append("-" * 70)
        lines.append(task.original_text)
        lines.append(f"\n润色后文本:")
        lines.append("-" * 70)
        lines.append(task.polished_text or "(未生成)")
        lines.append(f"\n问题详情 (共 {len(issues)} 个):")
        lines.append("-" * 70)
        
        for i, issue in enumerate(issues, 1):
            lines.append(f"\n{i}. 【{issue.issue_type.upper()}】{issue.severity.upper()}")
            lines.append(f"   原文: {issue.original_content}")
            lines.append(f"   建议: {issue.suggested_content}")
            lines.append(f"   原因: {issue.reason}")
            lines.append(f"   状态: {issue.status}")
        
        return "\n".join(lines)

    def _apply_suggestion(self, text: str, location: Dict[str, Any], replacement: str) -> str:
        """按位置应用建议内容"""
        try:
            start = int(location.get("start", 0))
            end = int(location.get("end", 0))
            if start < 0 or end <= start or end > len(text):
                return text
            return text[:start] + replacement + text[end:]
        except Exception:
            return text
    
    async def _polish_with_llm(self, original_text: str, polish_level: str = "standard") -> str:
        """
        使用Qwen LLM进行深度学术化改写
        
        Args:
            original_text: 原始文本
            polish_level: 润色级别 (light, standard, rigorous)
            
        Returns:
            学术化改写后的文本
        """
        if not settings.QWEN_API_KEY:
            logger.warning("Qwen API Key未配置，返回原文本")
            return original_text
        
        # 根据润色级别调整要求
        level_descriptions = {
            "light": "轻度学术化改写：主要提升表达的专业性和准确性",
            "standard": "中度学术化改写：同时提升学术规范性和表达严谨性",
            "rigorous": "深度学术化改写：全面提升学术水平，使用更多专业术语和学术表达"
        }
        
        level_desc = level_descriptions.get(polish_level, level_descriptions["standard"])
        
        prompt = (
            "你是资深的学术论文编辑，请对以下文本进行学术化改写。\n\n"
            "改写要求：\n"
            f"1. {level_desc}\n"
            "2. 转换为学术规范的表达方式：\n"
            "   - 使用学术化的词汇和短语\n"
            "   - 采用被动语态或名词化表达\n"
            "   - 添加必要的过渡词和逻辑连接词\n"
            "   - 优化句式结构使其更符合学术论文规范\n"
            "3. 保持原意不变，增强学术严谨性\n"
            "4. 改善时态和语态的一致性\n"
            "5. 不要改变文本的结构，只改进表达\n"
            "6. 不要添加任何额外的解释或注释\n\n"
            f"原文本：\n{original_text}\n\n"
            "请输出改写后的文本，确保明显高于原文的学术水平："
        )
        
        api_url = settings.QWEN_BASE_URL.rstrip("/") + "/chat/completions"
        headers = {
            "Authorization": f"Bearer {settings.QWEN_API_KEY}",
            "Content-Type": "application/json",
        }
        payload = {
            "model": settings.QWEN_MODEL_NAME,
            "messages": [
                {"role": "system", "content": "你是资深的学术论文编辑和写作顾问。"},
                {"role": "user", "content": prompt},
            ],
            "temperature": 0.3,
            "top_p": 0.9,
            "max_tokens": max(len(original_text) // 2 * 3, 500),  # 预留足够的token空间
            "stream": False,
        }
        
        try:
            async with aiohttp.ClientSession() as session:
                async with session.post(api_url, json=payload, headers=headers, timeout=60) as resp:
                    if resp.status >= 400:
                        error_text = await resp.text()
                        logger.error(f"Qwen API错误: status={resp.status}, response={error_text}")
                        logger.warning("Qwen API调用失败，返回原文本")
                        return original_text
                    
                    data = await resp.json()
            
            content = (
                data.get("choices", [{}])[0]
                .get("message", {})
                .get("content", "")
            ).strip()
            
            if not content:
                logger.warning("Qwen API返回空内容，返回原文本")
                return original_text
            
            logger.info("Qwen API学术化改写成功")
            return content
            
        except asyncio.TimeoutError:
            logger.error("Qwen API请求超时")
            return original_text
        except Exception as e:
            logger.error(f"调用Qwen API异常: {e}")
            return original_text





# ============================================================
# 多语言翻译服务
# ============================================================

class TranslationService:
    """多语言翻译服务"""

    def __init__(self, db: AsyncSession):
        self.db = db

    async def create_task(self, task_data, user_id: int) -> dict:
        """创建翻译任务"""
        logger.info("创建翻译任务")

        source_language = task_data.get("source_language") or "auto"
        target_language = task_data.get("target_language")
        input_text = task_data.get("input_text")
        domain = task_data.get("domain") or "general"

        if not target_language:
            raise ValueError("目标语言不能为空")
        if not input_text or not input_text.strip():
            raise ValueError("翻译文本不能为空")

        detected = source_language
        if source_language in (None, "", "auto"):
            detected = self._detect_language(input_text)

        terminology = await self._get_terminology_map(user_id, domain)
        translated_text = await self._translate_with_qwen(
            input_text.strip(),
            detected,
            target_language,
            terminology,
            domain
        )
        quality_score = self._estimate_quality(input_text, translated_text)

        task = TranslationTask(
            user_id=user_id,
            source_text=input_text.strip(),
            source_language=detected,
            target_language=target_language,
            translated_text=translated_text,
            status="completed",
            domain=domain,
            quality_score=quality_score,
        )
        self.db.add(task)
        await self.db.commit()
        await self.db.refresh(task)

        return self._format_task(task)

    async def list_tasks(self, user_id: int, skip: int, limit: int, status) -> dict:
        """获取任务列表"""
        query = select(TranslationTask).where(TranslationTask.user_id == user_id)
        if status:
            query = query.where(TranslationTask.status == status)
        query = query.order_by(TranslationTask.created_at.desc()).offset(skip).limit(limit)
        result = await self.db.execute(query)
        items = result.scalars().all()

        count_query = select(TranslationTask).where(TranslationTask.user_id == user_id)
        if status:
            count_query = count_query.where(TranslationTask.status == status)
        count_result = await self.db.execute(count_query)
        total = len(count_result.scalars().all())

        return {
            "total": total,
            "skip": skip,
            "limit": limit,
            "items": [self._format_task(task) for task in items]
        }

    async def get_task(self, task_id: str, user_id: int) -> dict:
        """获取任务详情"""
        query = select(TranslationTask).where(
            and_(TranslationTask.id == task_id, TranslationTask.user_id == user_id)
        )
        result = await self.db.execute(query)
        task = result.scalars().first()
        if not task:
            raise ValueError("任务不存在")
        return self._format_task(task)

    async def update_task(self, task_id: str, task_data, user_id: int) -> dict:
        """更新任务"""
        query = select(TranslationTask).where(
            and_(TranslationTask.id == task_id, TranslationTask.user_id == user_id)
        )
        result = await self.db.execute(query)
        task = result.scalars().first()
        if not task:
            raise ValueError("任务不存在")

        if task_data.get("input_text"):
            task.source_text = task_data["input_text"]
            task.status = "processing"
            terminology = await self._get_terminology_map(user_id, task.domain or "general")
            task.translated_text = await self._translate_with_qwen(
                task.source_text,
                task.source_language,
                task.target_language,
                terminology,
                task.domain or "general"
            )
            task.quality_score = self._estimate_quality(task.source_text, task.translated_text)
            task.status = "completed"

        task.updated_at = datetime.utcnow()
        await self.db.commit()
        await self.db.refresh(task)
        return self._format_task(task)

    async def delete_task(self, task_id: str, user_id: int) -> None:
        """删除任务"""
        query = select(TranslationTask).where(
            and_(TranslationTask.id == task_id, TranslationTask.user_id == user_id)
        )
        result = await self.db.execute(query)
        task = result.scalars().first()
        if not task:
            raise ValueError("任务不存在")
        await self.db.delete(task)
        await self.db.commit()

    async def get_terminology(self, user_id: int, domain: str = None) -> list:
        """获取术语库"""
        query = select(TranslationTerminology).where(TranslationTerminology.user_id == user_id)
        if domain:
            query = query.where(TranslationTerminology.domain == domain)
        result = await self.db.execute(query.order_by(TranslationTerminology.created_at.desc()))
        items = result.scalars().all()
        return [self._format_term(term) for term in items]

    async def add_terminology(self, term_data, user_id: int) -> dict:
        """添加术语"""
        original_term = term_data.get("original_term")
        translation = term_data.get("translation")
        domain = term_data.get("domain") or "general"
        if not original_term or not translation:
            raise ValueError("术语与译文不能为空")

        term = TranslationTerminology(
            user_id=user_id,
            original_term=original_term,
            translation=translation,
            domain=domain,
        )
        self.db.add(term)
        await self.db.commit()
        await self.db.refresh(term)
        return self._format_term(term)

    async def rate_translation(self, task_id: str, rating: int, feedback: str, user_id: int) -> dict:
        """评分翻译结果"""
        query = select(TranslationTask).where(
            and_(TranslationTask.id == task_id, TranslationTask.user_id == user_id)
        )
        result = await self.db.execute(query)
        task = result.scalars().first()
        if not task:
            raise ValueError("任务不存在")
        task.rating = rating
        task.feedback = feedback
        task.updated_at = datetime.utcnow()
        await self.db.commit()
        await self.db.refresh(task)
        return self._format_task(task)

    async def export_result(self, task_id: str, format_type: str, user_id: int) -> dict:
        """导出翻译结果"""
        query = select(TranslationTask).where(
            and_(TranslationTask.id == task_id, TranslationTask.user_id == user_id)
        )
        result = await self.db.execute(query)
        task = result.scalars().first()
        if not task:
            raise ValueError("任务不存在")

        if format_type == "json":
            return {
                "task": self._format_task(task),
                "export_time": datetime.utcnow().isoformat()
            }
        if format_type == "txt":
            lines = []
            lines.append("=" * 70)
            lines.append("多语言翻译导出")
            lines.append("=" * 70)
            lines.append(f"任务ID: {task.id}")
            lines.append(f"创建时间: {task.created_at}")
            lines.append(f"源语言: {task.source_language}")
            lines.append(f"目标语言: {task.target_language}")
            lines.append(f"领域: {task.domain}")
            lines.append("\n原文:")
            lines.append("-" * 70)
            lines.append(task.source_text)
            lines.append("\n译文:")
            lines.append("-" * 70)
            lines.append(task.translated_text or "")
            return "\n".join(lines)

        if format_type in {"pdf", "docx"}:
            file_url = self._export_as_document(task, format_type)
            return {
                "task_id": task.id,
                "format": format_type,
                "path": file_url
            }

        raise ValueError("不支持的导出格式")

    def _export_as_document(self, task: TranslationTask, format_type: str) -> str:
        upload_dir = settings.UPLOAD_DIR
        os.makedirs(upload_dir, exist_ok=True)

        filename = f"translation_{task.id}.{format_type}"
        file_path = os.path.join(upload_dir, filename)

        if format_type == "pdf":
            success = self._generate_translation_pdf(task, file_path)
        else:
            success = self._generate_translation_docx(task, file_path)

        if not success:
            raise ValueError(f"{format_type.upper()} 生成失败，请检查依赖是否安装")

        return f"/uploads/{filename}"

    def _generate_translation_pdf(self, task: TranslationTask, output_path: str) -> bool:
        try:
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.pdfbase import pdfmetrics
            from reportlab.pdfbase.ttfonts import TTFont
            from reportlab.pdfbase.cidfonts import UnicodeCIDFont
            import platform

            font_registered = False
            system = platform.system()
            font_name = "SimHei"
            font_paths: List[str] = []
            if system == "Windows":
                font_paths = [
                    "C:\\Windows\\Fonts\\simhei.ttf",
                    "C:\\Windows\\Fonts\\SimHei.ttf",
                    "C:\\Windows\\Fonts\\msyh.ttf",
                    "C:\\Windows\\Fonts\\msyhbd.ttf",
                ]
            elif system == "Darwin":
                font_paths = [
                    "/Library/Fonts/SimHei.ttf",
                    "/System/Library/Fonts/STHeiti Medium.ttc",
                    "/System/Library/Fonts/PingFang.ttc",
                ]
            else:
                font_paths = [
                    "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
                    "/usr/share/fonts/wqy-microhei/wqy-microhei.ttc",
                ]

            for font_path in font_paths:
                if os.path.exists(font_path):
                    try:
                        pdfmetrics.registerFont(TTFont(font_name, font_path))
                        font_registered = True
                        break
                    except Exception:
                        continue

            if not font_registered:
                try:
                    font_name = "STSong-Light"
                    pdfmetrics.registerFont(UnicodeCIDFont(font_name))
                    font_registered = True
                except Exception:
                    font_registered = False

            styles = getSampleStyleSheet()
            base_style = styles["BodyText"]
            if font_registered:
                base_style.fontName = font_name

            title_style = ParagraphStyle(
                "TranslationTitle",
                parent=styles["Title"],
                fontName=font_name if font_registered else styles["Title"].fontName,
            )
            heading_style = ParagraphStyle(
                "TranslationHeading",
                parent=styles["Heading2"],
                fontName=font_name if font_registered else styles["Heading2"].fontName,
            )

            doc = SimpleDocTemplate(output_path, pagesize=A4)
            elements = []
            elements.append(Paragraph("多语言翻译结果", title_style))
            elements.append(Spacer(1, 12))
            elements.append(Paragraph(f"源语言：{task.source_language}", base_style))
            elements.append(Paragraph(f"目标语言：{task.target_language}", base_style))
            elements.append(Paragraph(f"领域：{task.domain}", base_style))
            elements.append(Spacer(1, 12))

            elements.append(Paragraph("原文", heading_style))
            elements.append(Paragraph((task.source_text or "-").replace("\n", "<br/>") , base_style))
            elements.append(Spacer(1, 12))

            elements.append(Paragraph("译文", heading_style))
            elements.append(Paragraph((task.translated_text or "-").replace("\n", "<br/>") , base_style))

            doc.build(elements)
            logger.info(f"翻译 PDF 生成成功: {output_path}")
            return True
        except ImportError as e:
            logger.error(f"reportlab 未安装: {e}")
            return False
        except Exception as e:
            logger.error(f"翻译 PDF 生成失败: {type(e).__name__}: {e}", exc_info=True)
            return False

    def _generate_translation_docx(self, task: TranslationTask, output_path: str) -> bool:
        try:
            from docx import Document
            from docx.shared import Pt
            from docx.enum.text import WD_ALIGN_PARAGRAPH

            doc = Document()
            style = doc.styles["Normal"]
            style.font.name = "宋体"
            style.font.size = Pt(12)

            title_para = doc.add_heading("多语言翻译结果", level=1)
            title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

            doc.add_paragraph(f"源语言：{task.source_language}")
            doc.add_paragraph(f"目标语言：{task.target_language}")
            doc.add_paragraph(f"领域：{task.domain}")

            doc.add_heading("原文", level=2)
            for line in (task.source_text or "").splitlines():
                doc.add_paragraph(line)

            doc.add_heading("译文", level=2)
            for line in (task.translated_text or "").splitlines():
                doc.add_paragraph(line)

            doc.save(output_path)
            logger.info(f"翻译 Word 生成成功: {output_path}")
            return True
        except ImportError:
            logger.error("python-docx 未安装，请运行: pip install python-docx")
            return False
        except Exception as e:
            logger.error(f"翻译 Word 生成失败: {type(e).__name__}: {e}", exc_info=True)
            return False

    def _format_task(self, task: TranslationTask) -> dict:
        return {
            "id": task.id,
            "user_id": task.user_id,
            "source_language": task.source_language,
            "target_language": task.target_language,
            "input_text": task.source_text[:1000] + "..." if len(task.source_text) > 1000 else task.source_text,
            "translated_text": task.translated_text,
            "status": task.status,
            "domain": task.domain,
            "quality_score": task.quality_score,
            "created_at": task.created_at.isoformat() if task.created_at else None,
            "updated_at": task.updated_at.isoformat() if task.updated_at else None,
        }

    def _format_term(self, term: TranslationTerminology) -> dict:
        return {
            "id": term.id,
            "user_id": term.user_id,
            "original_term": term.original_term,
            "translation": term.translation,
            "domain": term.domain,
            "created_at": term.created_at.isoformat() if term.created_at else None,
        }

    def _detect_language(self, text: str) -> str:
        try:
            from langdetect import detect

            lang = detect(text)
            if lang in {"zh-cn", "zh-tw"}:
                return "zh"
            return lang
        except Exception:
            return "auto"

    async def _get_terminology_map(self, user_id: int, domain: str) -> dict:
        query = select(TranslationTerminology).where(
            and_(TranslationTerminology.user_id == user_id, TranslationTerminology.domain == domain)
        )
        result = await self.db.execute(query)
        items = result.scalars().all()
        return {term.original_term: term.translation for term in items}

    async def _translate_with_qwen(
        self,
        text: str,
        source_language: str,
        target_language: str,
        terminology: dict,
        domain: str,
    ) -> str:
        if not settings.QWEN_API_KEY:
            return "（未配置 QWEN_API_KEY，无法进行翻译。请在后端 .env 中配置后重试。）"

        terminology_hint = "\n".join(
            [f"- {k} => {v}" for k, v in terminology.items()]
        )
        if terminology_hint:
            terminology_hint = "术语对照表：\n" + terminology_hint
        else:
            terminology_hint = "术语对照表：无"

        prompt = (
            "你是专业翻译助手，请按要求输出最终译文：\n"
            "1) 只输出译文，不要解释\n"
            "2) 保持语气自然、专业、简洁\n"
            "3) 优先使用给定术语对照表\n\n"
            f"翻译领域：{domain}\n"
            f"源语言：{source_language}\n"
            f"目标语言：{target_language}\n"
            f"{terminology_hint}\n\n"
            f"原文：\n{text}"
        )

        api_url = settings.QWEN_BASE_URL.rstrip("/") + "/chat/completions"
        headers = {
            "Authorization": f"Bearer {settings.QWEN_API_KEY}",
            "Content-Type": "application/json",
        }
        payload = {
            "model": settings.QWEN_MODEL_NAME,
            "messages": [
                {"role": "system", "content": "你是专业的多语言翻译助手。"},
                {"role": "user", "content": prompt},
            ],
            "temperature": 0.2,
            "top_p": 0.9,
            "max_tokens": 2048,
            "stream": False,
        }

        async with aiohttp.ClientSession() as session:
            async with session.post(api_url, json=payload, headers=headers, timeout=60) as resp:
                if resp.status >= 400:
                    raise ValueError("翻译失败")
                data = await resp.json()
        content = (
            data.get("choices", [{}])[0]
            .get("message", {})
            .get("content", "")
        )
        return content.strip()

    def _estimate_quality(self, source: str, translated: str) -> float:
        if not source or not translated:
            return 0.0
        ratio = min(len(translated) / max(len(source), 1), 1.0)
        return round(max(0.2, 1 - abs(ratio - 0.9)), 3)


# ============================================================
# PPT生成服务
# ============================================================

class PPTService:
    """PPT生成服务"""
    
    def __init__(self, db: AsyncSession):
        self.db = db
    
    async def create_project(self, project_data, user_id: int) -> dict:
        """创建PPT项目"""
        logger.info("创建PPT项目")

        title = project_data.get("title")
        content = project_data.get("source_content")
        description = project_data.get("description")
        theme = project_data.get("theme") or "classic"
        theme_palette = project_data.get("theme_palette")

        if not title:
            raise ValueError("项目名称不能为空")
        if not content or not content.strip():
            raise ValueError("请输入内容")

        project = PPTProject(
            user_id=user_id,
            title=title,
            description=description,
            content=content.strip(),
            theme=theme,
            theme_palette=json.dumps(theme_palette, ensure_ascii=False) if isinstance(theme_palette, dict) else None,
            status="draft",
        )
        self.db.add(project)
        await self.db.commit()
        await self.db.refresh(project)
        return self._format_project(project)

    async def create_project_from_file(
        self,
        title: str,
        file,
        description: Optional[str],
        theme: Optional[str],
        theme_palette: Optional[str],
        user_id: int
    ) -> dict:
        """从文件导入创建PPT项目"""
        palette_value = None
        if theme_palette:
            try:
                palette_value = json.loads(theme_palette)
            except Exception:
                palette_value = None
        filename = file.filename or "document"
        ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
        if ext not in {"pdf", "docx", "md", "markdown"}:
            raise ValueError("仅支持 PDF/DOCX/Markdown 文件")

        content = await file.read()
        if not content:
            raise ValueError("文件内容为空")

        os.makedirs(settings.UPLOAD_DIR, exist_ok=True)
        saved_name = f"ppt_source_{uuid.uuid4().hex}.{ext}"
        saved_path = os.path.join(settings.UPLOAD_DIR, saved_name)
        with open(saved_path, "wb") as f:
            f.write(content)

        parsed_text = self._parse_ppt_source_file(saved_path, ext)
        if not parsed_text.strip():
            raise ValueError("未解析到有效内容")

        project = PPTProject(
            user_id=user_id,
            title=title,
            description=description,
            content=parsed_text.strip(),
            theme=theme or "classic",
            theme_palette=json.dumps(palette_value, ensure_ascii=False) if isinstance(palette_value, dict) else None,
            status="draft",
        )
        self.db.add(project)
        await self.db.commit()
        await self.db.refresh(project)
        return self._format_project(project)

    async def list_projects(self, user_id: int, skip: int, limit: int, status) -> dict:
        """获取项目列表"""
        query = select(PPTProject).where(PPTProject.user_id == user_id)
        if status:
            query = query.where(PPTProject.status == status)
        query = query.order_by(PPTProject.created_at.desc()).offset(skip).limit(limit)
        result = await self.db.execute(query)
        items = result.scalars().all()

        count_query = select(PPTProject).where(PPTProject.user_id == user_id)
        if status:
            count_query = count_query.where(PPTProject.status == status)
        count_result = await self.db.execute(count_query)
        total = len(count_result.scalars().all())

        return {
            "total": total,
            "skip": skip,
            "limit": limit,
            "items": [self._format_project(item) for item in items],
        }

    async def get_project(self, project_id: str, user_id: int) -> dict:
        """获取项目详情"""
        query = select(PPTProject).where(
            and_(PPTProject.id == project_id, PPTProject.user_id == user_id)
        )
        result = await self.db.execute(query)
        project = result.scalars().first()
        if not project:
            raise ValueError("项目不存在")
        return self._format_project(project, include_slides=True)

    async def update_project(self, project_id: str, project_data, user_id: int) -> dict:
        """更新项目"""
        query = select(PPTProject).where(
            and_(PPTProject.id == project_id, PPTProject.user_id == user_id)
        )
        result = await self.db.execute(query)
        project = result.scalars().first()
        if not project:
            raise ValueError("项目不存在")

        if project_data.get("title"):
            project.title = project_data["title"]
        if project_data.get("description") is not None:
            project.description = project_data.get("description")
        if project_data.get("theme"):
            project.theme = project_data.get("theme")
        if project_data.get("theme_palette") is not None:
            theme_palette = project_data.get("theme_palette")
            project.theme_palette = (
                json.dumps(theme_palette, ensure_ascii=False)
                if isinstance(theme_palette, dict)
                else None
            )
        project.updated_at = datetime.utcnow()
        await self.db.commit()
        await self.db.refresh(project)
        return self._format_project(project)

    async def delete_project(self, project_id: str, user_id: int) -> None:
        """删除项目"""
        query = select(PPTProject).where(
            and_(PPTProject.id == project_id, PPTProject.user_id == user_id)
        )
        result = await self.db.execute(query)
        project = result.scalars().first()
        if not project:
            raise ValueError("项目不存在")
        await self.db.delete(project)
        await self.db.commit()

    async def generate_slides(
        self,
        project_id: str,
        user_id: int,
        tone: str,
        theme: Optional[str],
        theme_palette: Optional[dict]
    ) -> dict:
        """生成幻灯片"""
        logger.info(f"开始生成PPT幻灯片: project_id={project_id}, tone={tone}")
        query = select(PPTProject).where(
            and_(PPTProject.id == project_id, PPTProject.user_id == user_id)
        )
        result = await self.db.execute(query)
        project = result.scalars().first()
        if not project:
            raise ValueError("项目不存在")

        logger.info(f"找到项目: title={project.title}, content_length={len(project.content or '')}")

        if theme:
            project.theme = theme
        if theme_palette is not None:
            project.theme_palette = (
                json.dumps(theme_palette, ensure_ascii=False)
                if isinstance(theme_palette, dict)
                else None
            )

        logger.info("开始调用Qwen API生成大纲...")
        outline = await self._build_outline(project.title, project.content, tone)
        logger.info(f"大纲生成成功，幻灯片数量: {len(outline.get('slides', []))}")
        
        slides = self._outline_to_slides(outline)
        logger.info(f"幻灯片转换完成，最终数量: {len(slides)}")

        project.outline_json = json.dumps(outline, ensure_ascii=False)
        project.slides_json = json.dumps(slides, ensure_ascii=False)
        project.status = "completed"
        project.updated_at = datetime.utcnow()
        await self.db.commit()
        await self.db.refresh(project)
        logger.info(f"PPT幻灯片生成完成: project_id={project_id}")
        return self._format_project(project, include_slides=True)

    async def get_slides(self, project_id: str, user_id: int) -> list:
        """获取幻灯片列表"""
        query = select(PPTProject).where(
            and_(PPTProject.id == project_id, PPTProject.user_id == user_id)
        )
        result = await self.db.execute(query)
        project = result.scalars().first()
        if not project:
            raise ValueError("项目不存在")
        slides = self._safe_json(project.slides_json, default=[])
        return slides

    async def export_pptx(self, project_id: str, user_id: int, format_type: str = "pptx") -> dict:
        """导出PPTX文件"""
        if format_type != "pptx":
            raise ValueError("暂不支持该格式")

        query = select(PPTProject).where(
            and_(PPTProject.id == project_id, PPTProject.user_id == user_id)
        )
        result = await self.db.execute(query)
        project = result.scalars().first()
        if not project:
            raise ValueError("项目不存在")

        slides = self._safe_json(project.slides_json, default=[])
        if not slides:
            raise ValueError("请先生成幻灯片")

        expanded_slides = await self._expand_slides_for_export(project.title, slides)

        file_url = self._export_as_pptx(project, expanded_slides)
        project.file_path = file_url
        project.updated_at = datetime.utcnow()
        await self.db.commit()
        await self.db.refresh(project)
        return {"project_id": project.id, "format": "pptx", "path": file_url}

    async def _expand_slides_for_export(self, title: str, slides: list) -> list:
        if not llm_service.check_availability():
            return slides

        async def _expand_slide(slide: dict) -> dict:
            bullets = slide.get("bullets") or []
            if not bullets:
                return slide
            expanded = await llm_service.expand_ppt_bullets(
                title=title,
                slide_title=slide.get("title") or "",
                bullets=bullets,
            )
            return {
                **slide,
                "bullets": expanded,
            }

        results = await asyncio.gather(*[_expand_slide(slide) for slide in slides])
        return results

    async def _build_outline(self, title: str, content: str, tone: str) -> dict:
        logger.info(f"_build_outline: QWEN_API_KEY={'已配置' if settings.QWEN_API_KEY else '未配置'}")
        if not settings.QWEN_API_KEY:
            logger.warning("Qwen API Key未配置，使用fallback生成大纲")
            return self._fallback_outline(title, content)

        prompt = (
            "你是专业PPT策划助手，请输出严格JSON格式，不要解释。\n"
            "输出格式：{\"title\": str, \"slides\": [{\"title\": str, \"bullets\": [str], \"notes\": str}]}\n"
            "要求：1) 幻灯片6-10页 2) 重点清晰 3) 语气自然专业 4) 不要出现AI/模型等字眼\n"
            f"演示主题：{title}\n"
            f"表达风格：{tone}\n"
            f"输入内容：\n{content[:60000]}"
        )

        api_url = settings.QWEN_BASE_URL.rstrip("/") + "/chat/completions"
        logger.info(f"调用Qwen API: {api_url}")
        
        headers = {
            "Authorization": f"Bearer {settings.QWEN_API_KEY}",
            "Content-Type": "application/json",
        }
        payload = {
            "model": settings.QWEN_MODEL_NAME,
            "messages": [
                {"role": "system", "content": "你是专业PPT策划助手。"},
                {"role": "user", "content": prompt},
            ],
            "temperature": 0.2,
            "top_p": 0.9,
            "max_tokens": 1800,
            "stream": False,
        }

        try:
            async with aiohttp.ClientSession() as session:
                async with session.post(api_url, json=payload, headers=headers, timeout=60) as resp:
                    logger.info(f"Qwen API响应状态: {resp.status}")
                    if resp.status >= 400:
                        error_text = await resp.text()
                        logger.error(f"Qwen API错误: status={resp.status}, response={error_text}")
                        raise ValueError(f"生成失败: API返回{resp.status}")
                    data = await resp.json()
                    logger.info(f"Qwen API响应成功")
        except aiohttp.ClientError as e:
            logger.error(f"Qwen API请求异常: {e}")
            logger.warning("使用fallback生成大纲")
            return self._fallback_outline(title, content)
        except Exception as e:
            logger.error(f"调用Qwen API异常: {e}")
            logger.warning("使用fallback生成大纲")
            return self._fallback_outline(title, content)
            
        content = (
            data.get("choices", [{}])[0]
            .get("message", {})
            .get("content", "")
        )
        logger.info(f"Qwen返回内容长度: {len(content)}")
        
        try:
            outline = json.loads(content)
            logger.info("成功解析JSON大纲")
            return outline
        except Exception as e:
            logger.error(f"JSON解析失败: {e}, content={content[:200]}")
            logger.warning("使用fallback生成大纲")
            return self._fallback_outline(title, content)

    def _fallback_outline(self, title: str, content: str) -> dict:
        paragraphs = [p.strip() for p in content.split("\n") if p.strip()][:6]
        slides = []
        for idx, p in enumerate(paragraphs, 1):
            slides.append({
                "title": f"要点 {idx}",
                "bullets": [p[:80]],
                "notes": ""
            })
        return {"title": title, "slides": slides}

    def _outline_to_slides(self, outline: dict) -> list:
        slides = outline.get("slides") or []
        return [
            {
                "title": slide.get("title") or "未命名",
                "bullets": slide.get("bullets") or [],
                "notes": slide.get("notes") or "",
            }
            for slide in slides
        ]

    def _export_as_pptx(self, project: PPTProject, slides: list) -> str:
        from pptx import Presentation
        from pptx.util import Pt
        from pptx.dml.color import RGBColor

        upload_dir = settings.UPLOAD_DIR
        os.makedirs(upload_dir, exist_ok=True)

        filename = f"ppt_{project.id}.pptx"
        file_path = os.path.join(upload_dir, filename)

        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = project.title
        if project.description:
            slide.placeholders[1].text = project.description
        self._apply_theme(slide, project.theme, project.theme_palette, RGBColor)

        bullet_layout = prs.slide_layouts[1]
        for item in slides:
            slide = prs.slides.add_slide(bullet_layout)
            slide.shapes.title.text = item.get("title") or ""
            body = slide.shapes.placeholders[1].text_frame
            body.clear()
            for bullet in item.get("bullets", []):
                p = body.add_paragraph()
                p.text = bullet
                p.level = 0
                p.font.size = Pt(20)
            self._apply_theme(slide, project.theme, project.theme_palette, RGBColor)

        prs.save(file_path)
        return f"/uploads/{filename}"

    def _parse_ppt_source_file(self, path: str, ext: str) -> str:
        if ext in {"md", "markdown"}:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        if ext == "docx":
            from docx import Document

            doc = Document(path)
            return "\n".join([p.text for p in doc.paragraphs if p.text])
        if ext == "pdf":
            import pdfplumber

            texts = []
            with pdfplumber.open(path) as pdf:
                for page in pdf.pages:
                    texts.append(page.extract_text() or "")
            return "\n".join(texts)
        return ""

    def _format_project(self, project: PPTProject, include_slides: bool = False) -> dict:
        payload = {
            "id": project.id,
            "title": project.title,
            "description": project.description,
            "status": project.status,
            "created_at": project.created_at.isoformat() if project.created_at else None,
            "updated_at": project.updated_at.isoformat() if project.updated_at else None,
            "file_path": project.file_path,
            "theme": project.theme,
        }
        if include_slides:
            payload["slides"] = self._safe_json(project.slides_json, default=[])
            payload["outline"] = self._safe_json(project.outline_json, default={})
            payload["theme_palette"] = self._safe_json(project.theme_palette, default=None)
        return payload

    def _safe_json(self, value: Optional[str], default):
        if not value:
            return default
        try:
            return json.loads(value)
        except Exception:
            return default

    def _apply_theme(self, slide, theme: Optional[str], theme_palette: Optional[str], RGBColor) -> None:
        theme_key = (theme or "classic").lower()
        themes = {
            "classic": {"bg": RGBColor(255, 255, 255), "text": RGBColor(30, 41, 59)},
            "dark": {"bg": RGBColor(15, 23, 42), "text": RGBColor(248, 250, 252)},
            "ocean": {"bg": RGBColor(226, 240, 255), "text": RGBColor(15, 23, 42)},
            "forest": {"bg": RGBColor(232, 245, 233), "text": RGBColor(27, 94, 32)},
        }
        palette = themes.get(theme_key, themes["classic"])
        custom_palette = self._safe_json(theme_palette, default=None)
        if isinstance(custom_palette, dict):
            bg = custom_palette.get("bg")
            text = custom_palette.get("text")
            if isinstance(bg, str) and bg.startswith("#"):
                palette = {
                    "bg": self._hex_to_rgb(bg, RGBColor),
                    "text": self._hex_to_rgb(text, RGBColor) if isinstance(text, str) else palette["text"],
                }

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = palette["bg"]

        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = palette["text"]

    def _hex_to_rgb(self, value: str, RGBColor):
        hex_value = value.lstrip("#")
        if len(hex_value) != 6:
            return RGBColor(255, 255, 255)
        r = int(hex_value[0:2], 16)
        g = int(hex_value[2:4], 16)
        b = int(hex_value[4:6], 16)
        return RGBColor(r, g, b)


# ============================================================
# 周报生成服务 - 已在 report_service.py 中完整实现
# ============================================================

# 为了保持向后兼容性，在此导入并代理实现
from app.services.report_service import WeeklyReportService as ReportService

